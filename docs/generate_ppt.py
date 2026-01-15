#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
ERP系统优势分析PPT生成脚本
运行方式: python generate_ppt.py
输出文件: output/ERP系统优势分析.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.dml.color import RGBColor
import os

# 颜色定义
PRIMARY_COLOR = RGBColor(0x1E, 0x88, 0xE5)  # 主色调蓝色
SECONDARY_COLOR = RGBColor(0x42, 0x42, 0x42)  # 深灰色
ACCENT_COLOR = RGBColor(0x43, 0xA0, 0x47)  # 绿色强调
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)


def add_title_slide(prs, title, subtitle=""):
    """添加标题幻灯片"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景色
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = PRIMARY_COLOR
    background.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5), Inches(9), Inches(1.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.2), Inches(9), Inches(0.8)
        )
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    
    return slide


def add_section_slide(prs, title):
    """添加章节标题幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 左侧色块
    left_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.3), prs.slide_height
    )
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = PRIMARY_COLOR
    left_bar.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(2.8), Inches(8.5), Inches(1.2)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = SECONDARY_COLOR
    
    return slide


def add_content_slide(prs, title, bullet_points, subtitle=""):
    """添加内容幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 顶部色条
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.1)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = PRIMARY_COLOR
    top_bar.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = SECONDARY_COLOR
    
    # 副标题
    start_y = 1.2
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.1), Inches(9), Inches(0.5)
        )
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0x75, 0x75, 0x75)
        start_y = 1.6
    
    # 要点列表
    content_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(start_y), Inches(9), Inches(5)
    )
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + point
        p.font.size = Pt(20)
        p.font.color.rgb = SECONDARY_COLOR
        p.space_after = Pt(12)
    
    return slide


def add_comparison_slide(prs, title, headers, rows):
    """添加对比表格幻灯片"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 顶部色条
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.1)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = PRIMARY_COLOR
    top_bar.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.7)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = SECONDARY_COLOR
    
    # 表格
    cols = len(headers)
    table_rows = len(rows) + 1
    
    x, y, cx, cy = Inches(0.3), Inches(1.2), Inches(9.4), Inches(0.5 * table_rows)
    table = slide.shapes.add_table(table_rows, cols, x, y, cx, cy).table
    
    # 设置列宽
    col_width = Inches(9.4 / cols)
    for i in range(cols):
        table.columns[i].width = col_width
    
    # 表头
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY_COLOR
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    
    # 数据行
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = SECONDARY_COLOR
            p.alignment = PP_ALIGN.CENTER
            
            # 交替行背景
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
    
    return slide


def add_highlight_slide(prs, title, points):
    """添加亮点展示幻灯片（带图标风格）"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 顶部色条
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.1)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = PRIMARY_COLOR
    top_bar.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.7)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = SECONDARY_COLOR
    
    # 亮点卡片
    cols = min(3, len(points))
    card_width = 2.8
    card_height = 2.5
    start_x = (10 - cols * card_width - (cols - 1) * 0.3) / 2
    
    for i, (icon, point_title, desc) in enumerate(points):
        row = i // 3
        col = i % 3
        x = start_x + col * (card_width + 0.3)
        y = 1.3 + row * (card_height + 0.3)
        
        # 卡片背景
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(card_width), Inches(card_height)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GRAY
        card.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        
        # 图标文字
        icon_box = slide.shapes.add_textbox(
            Inches(x), Inches(y + 0.2), Inches(card_width), Inches(0.6)
        )
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(32)
        p.alignment = PP_ALIGN.CENTER
        
        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(x + 0.1), Inches(y + 0.9), Inches(card_width - 0.2), Inches(0.5)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = point_title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = SECONDARY_COLOR
        p.alignment = PP_ALIGN.CENTER
        
        # 描述
        desc_box = slide.shapes.add_textbox(
            Inches(x + 0.1), Inches(y + 1.4), Inches(card_width - 0.2), Inches(1)
        )
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
        p.alignment = PP_ALIGN.CENTER
    
    return slide


def create_ppt():
    """创建PPT"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # 封面
    add_title_slide(prs, "ERP系统优势分析", "面向非标/项目型制造企业的一体化管理平台")
    
    # 目录
    add_content_slide(prs, "目录", [
        "系统定位与价值主张",
        "核心功能优势",
        "技术架构优势",
        "成本对比分析",
        "竞品对比",
        "适用客户画像",
        "总结"
    ])
    
    # ===== 第一章：系统定位 =====
    add_section_slide(prs, "01 系统定位与价值主张")
    
    add_content_slide(prs, "市场痛点", [
        "通用ERP面向批量生产企业，无法适应'一单一设计'的非标模式",
        "CRM、PLM、ERP、MES各自独立，数据孤岛严重",
        "SAP/Oracle等国际品牌价格昂贵，实施周期长",
        "国产ERP定制能力弱，无法满足个性化需求",
        "项目成本核算困难，物料采购无法追溯到具体项目"
    ], subtitle="非标/项目型制造企业面临的信息化困境")
    
    add_content_slide(prs, "我们的解决方案", [
        "专为非标/项目型制造企业设计，每个项目独立BOM",
        "CRM+PLM+ERP+MES+OA五大系统一体化，数据天然打通",
        "销售订单→项目→BOM→采购→生产→发货，全程可追溯",
        "现代化技术架构，私有化部署，数据安全可控",
        "开源可定制，成本可降低80%以上"
    ], subtitle="一套系统解决所有问题")
    
    # ===== 第二章：功能优势 =====
    add_section_slide(prs, "02 核心功能优势")
    
    add_highlight_slide(prs, "六大核心优势", [
        ("🎯", "项目型管理", "专为非标企业设计\n每个项目独立BOM"),
        ("🔗", "全流程打通", "销售→项目→采购→生产\n数据自动流转"),
        ("📊", "状态实时追踪", "BOM采购状态全程可视\n10种状态精准跟踪"),
        ("📁", "一体化平台", "CRM+PLM+ERP+MES+OA\n无需多系统集成"),
        ("⚡", "简化流程", "去掉冗余环节\n适合中小企业"),
        ("🔧", "灵活定制", "源码可控\n满足个性化需求")
    ])
    
    add_content_slide(prs, "销售-项目-采购 深度打通", [
        "创建项目时可关联销售订单，自动带出客户、金额、交期",
        "项目BOM自动继承，支持按项目筛选和导出询价",
        "采购申请可批量导入，自动匹配项目BOM",
        "BOM采购状态实时更新：未下单→申请中→已批准→已下单→在途→已收货",
        "已关联项目的订单自动从选择列表中移除，防止重复"
    ], subtitle="数据自动流转，无需手工录入")
    
    add_content_slide(prs, "简化的采购询价流程", [
        "传统ERP：采购申请→询价单→比价→采购订单（流程复杂）",
        "我们的方案：采购申请页面直接导出询价BOM",
        "支持按'有图/无图'、'物料类型'、'版本/品牌'筛选",
        "询价在系统外进行（邮件/微信/电话），更灵活高效",
        "适合非标企业的询价特点：需要和供应商反复沟通技术细节"
    ], subtitle="去掉不必要的环节，聚焦核心业务")
    
    # ===== 第三章：技术优势 =====
    add_section_slide(prs, "03 技术架构优势")
    
    add_comparison_slide(prs, "技术栈对比", 
        ["组件", "传统ERP", "我们的系统"],
        [
            ["前端", "Java Swing/JSP", "Vue 3 + Element Plus"],
            ["后端", "Java/ABAP", "Python Django REST"],
            ["数据库", "Oracle/SQL Server", "PostgreSQL（免费）"],
            ["消息队列", "IBM MQ", "Redis + Celery"],
            ["搜索引擎", "无/商业", "Elasticsearch"],
            ["部署方式", "复杂、依赖多", "Docker一键部署"],
        ]
    )
    
    add_content_slide(prs, "技术优势详解", [
        "前端体验：单页应用，操作流畅，无需频繁刷新页面",
        "私有化部署：数据完全自主掌控，安全合规",
        "RESTful API：方便与其他系统集成",
        "Webhook支持：可对接企业微信、钉钉等",
        "源代码可控：可根据业务需求深度定制",
        "二次开发效率高：Python/Vue生态丰富"
    ])
    
    # ===== 第四章：成本对比 =====
    add_section_slide(prs, "04 成本对比分析")
    
    add_comparison_slide(prs, "成本对比（5年TCO）",
        ["成本项", "SAP/Oracle", "用友/金蝶", "我们的系统"],
        [
            ["软件许可", "100-500万", "20-80万", "0（开源）"],
            ["数据库", "50-100万", "5-20万", "0（PostgreSQL）"],
            ["实施费用", "100-300万", "10-50万", "可自主实施"],
            ["年维护费", "15-20%许可费", "10-15%许可费", "0"],
            ["二次开发", "昂贵", "较贵", "成本可控"],
            ["总计（5年）", "300-1000万", "50-200万", "10-50万"],
        ]
    )
    
    add_content_slide(prs, "成本优势总结", [
        "软件许可费：基于开源技术，无许可费用",
        "数据库费用：使用PostgreSQL，无需Oracle授权",
        "实施费用：系统简洁，可自主实施或低成本外包",
        "维护费用：无强制年度维护费",
        "二次开发：Python/Vue开发效率高，人才市场充裕",
        "总体拥有成本（TCO）可降低80%以上"
    ], subtitle="让中小企业也能用得起专业ERP")
    
    # ===== 第五章：竞品对比 =====
    add_section_slide(prs, "05 竞品对比")
    
    add_comparison_slide(prs, "功能对比",
        ["功能", "SAP B1", "用友U8+", "金蝶云星空", "我们"],
        [
            ["项目型制造", "一般", "一般", "较好", "⭐⭐⭐⭐⭐"],
            ["PLM集成", "单独购买", "单独购买", "部分", "⭐⭐⭐⭐⭐原生"],
            ["MES集成", "单独购买", "单独购买", "单独购买", "⭐⭐⭐⭐原生"],
            ["实施难度", "高", "中", "中", "⭐⭐⭐⭐⭐低"],
            ["定制灵活", "低", "中", "中", "⭐⭐⭐⭐⭐高"],
            ["总成本", "高", "中", "中高", "⭐⭐⭐⭐⭐低"],
        ]
    )
    
    # ===== 第六章：客户画像 =====
    add_section_slide(prs, "06 适用客户画像")
    
    add_content_slide(prs, "最适合的企业类型", [
        "非标设备制造：自动化设备、专用机械、检测设备等",
        "项目型生产：一单一设计、按订单生产（MTO/ETO）",
        "中小型制造企业：50-500人规模",
        "成长型企业：需要灵活、可扩展的信息化系统",
        "注重数据安全：需要私有化部署的企业",
        "有IT能力：希望自主掌控系统的企业"
    ])
    
    add_content_slide(prs, "典型应用场景", [
        "接到客户订单后，创建项目并关联销售订单",
        "技术部门设计BOM，按项目管理物料清单",
        "采购部门根据BOM询价、下单，全程可追溯",
        "生产部门按项目领料、生产、质检",
        "财务部门按项目归集成本、核算利润",
        "管理层实时查看项目进度、成本、利润"
    ], subtitle="从接单到交付的全流程数字化管理")
    
    # ===== 总结 =====
    add_section_slide(prs, "07 总结")
    
    add_highlight_slide(prs, "核心卖点", [
        ("🎯", "专业定位", "专为非标/项目型\n制造企业设计"),
        ("🔗", "一体化", "CRM+PLM+ERP+MES+OA\n五位一体"),
        ("💰", "高性价比", "TCO降低80%\n中小企业可承受"),
        ("🔒", "数据安全", "私有化部署\n完全自主可控"),
        ("⚡", "快速上线", "Docker部署\n2周可上线"),
        ("🔧", "灵活定制", "源码可控\n满足个性需求")
    ])
    
    # 封底
    add_title_slide(prs, "谢谢！", "如有疑问，欢迎交流")
    
    # 保存
    output_dir = os.path.dirname(os.path.abspath(__file__))
    output_path = os.path.join(output_dir, "output", "ERP系统优势分析.pptx")
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"✅ PPT已生成: {output_path}")
    return output_path


if __name__ == "__main__":
    create_ppt()
